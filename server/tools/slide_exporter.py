#!/usr/bin/env python3
"""
Slide Exporter Tool
Exports PowerPoint presentations to markdown format with PNG images
Uses LibreOffice for slide rendering
"""

import os
import sys
import json
import subprocess
import shutil
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional

from pptx import Presentation


class SlideExporter:
    """Export PowerPoint presentations to various formats"""
    
    def __init__(self, presentations_dir: str = "./presentations"):
        self.presentations_dir = Path(presentations_dir)
        self.exports_dir = self.presentations_dir / "exports"
        self.exports_dir.mkdir(parents=True, exist_ok=True)
    
    def export_to_markdown(self, pptx_file: Path) -> Path:
        """
        Export a PowerPoint presentation to markdown with PNG images
        
        Args:
            pptx_file: Path to the PPTX file
        
        Returns:
            Path to the generated markdown file
        """
        if not pptx_file.exists():
            raise FileNotFoundError(f"Presentation file not found: {pptx_file}")
        
        # Create export directory for this presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        export_name = f"{pptx_file.stem}_{timestamp}"
        export_dir = self.exports_dir / export_name
        export_dir.mkdir(parents=True, exist_ok=True)
        
        # Copy PPTX to export directory
        export_pptx = export_dir / pptx_file.name
        shutil.copy2(pptx_file, export_pptx)
        
        # Load presentation
        prs = Presentation(pptx_file)
        
        # Generate PNG images using LibreOffice
        images_dir = export_dir / "images"
        images_dir.mkdir(exist_ok=True)
        self._generate_slide_images(export_pptx, images_dir)
        
        # Create markdown content
        markdown_content = self._create_markdown(prs, export_name, images_dir)
        
        # Save markdown file
        markdown_file = export_dir / f"{export_name}_overview.md"
        markdown_file.write_text(markdown_content, encoding='utf-8')
        
        # Generate PDF if requested
        self._generate_pdf(export_pptx, export_dir)
        
        print(f"✅ Export completed: {export_dir}")
        print(f"📄 Markdown: {markdown_file}")
        
        return markdown_file
    
    def _generate_slide_images(self, pptx_file: Path, output_dir: Path) -> List[Path]:
        """
        Generate PNG images from PowerPoint slides using LibreOffice
        
        Args:
            pptx_file: Path to PPTX file
            output_dir: Directory to save images
        
        Returns:
            List of generated image paths
        """
        print(f"🖼️  Generating slide images...")
        
        # First convert to PDF
        pdf_cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir.parent),
            str(pptx_file)
        ]
        
        try:
            result = subprocess.run(pdf_cmd, capture_output=True, text=True, check=True)
            pdf_file = output_dir.parent / f"{pptx_file.stem}.pdf"
            
            if pdf_file.exists():
                # Convert PDF to PNG images using ImageMagick or pdftoppm
                if shutil.which("pdftoppm"):
                    self._pdf_to_png_pdftoppm(pdf_file, output_dir)
                elif shutil.which("convert"):
                    self._pdf_to_png_imagemagick(pdf_file, output_dir)
                else:
                    # Fallback: Direct PPTX to PNG conversion
                    self._pptx_to_png_direct(pptx_file, output_dir)
                
                # Clean up PDF
                pdf_file.unlink()
            
        except subprocess.CalledProcessError as e:
            print(f"⚠️  Warning: Could not generate images via LibreOffice: {e}")
            # Try direct conversion
            self._pptx_to_png_direct(pptx_file, output_dir)
        
        # Return list of generated images
        images = sorted(output_dir.glob("*.png"))
        print(f"✅ Generated {len(images)} slide images")
        return images
    
    def _pdf_to_png_pdftoppm(self, pdf_file: Path, output_dir: Path):
        """Convert PDF to PNG using pdftoppm"""
        cmd = [
            "pdftoppm",
            "-png",
            "-r", "150",  # DPI
            str(pdf_file),
            str(output_dir / "slide")
        ]
        subprocess.run(cmd, check=True)
        
        # Rename files to have consistent naming
        for i, img in enumerate(sorted(output_dir.glob("slide-*.png")), 1):
            img.rename(output_dir / f"slide-{i:03d}.png")
    
    def _pdf_to_png_imagemagick(self, pdf_file: Path, output_dir: Path):
        """Convert PDF to PNG using ImageMagick"""
        cmd = [
            "convert",
            "-density", "150",
            str(pdf_file),
            "-quality", "90",
            str(output_dir / "slide-%03d.png")
        ]
        subprocess.run(cmd, check=True)
    
    def _pptx_to_png_direct(self, pptx_file: Path, output_dir: Path):
        """Direct PPTX to PNG conversion using LibreOffice"""
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "png",
            "--outdir",
            str(output_dir),
            str(pptx_file)
        ]
        
        try:
            subprocess.run(cmd, check=True)
            # Rename the single image if only one was created
            single_image = output_dir / f"{pptx_file.stem}.png"
            if single_image.exists():
                single_image.rename(output_dir / "slide-001.png")
        except subprocess.CalledProcessError:
            print("⚠️  Could not generate images. LibreOffice may not be installed.")
    
    def _create_markdown(self, prs: Presentation, export_name: str, images_dir: Path) -> str:
        """
        Create markdown documentation for the presentation
        
        Args:
            prs: PowerPoint presentation object
            export_name: Name of the export
            images_dir: Directory containing slide images
        
        Returns:
            Markdown content as string
        """
        images = sorted(images_dir.glob("*.png"))
        
        # Start markdown content
        md_lines = [
            f"# {export_name.replace('_', ' ').title()}",
            "",
            "## Metadata",
            f"- **Created**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            f"- **Total Slides**: {len(prs.slides)}",
            f"- **Estimated Duration**: {len(prs.slides) * 2} minutes",
            f"- **Generated By**: Claude Code Slide Agent",
            "",
            "## Table of Contents",
            ""
        ]
        
        # Add TOC
        for i, slide in enumerate(prs.slides, 1):
            title = self._get_slide_title(slide) or f"Slide {i}"
            md_lines.append(f"{i}. [{title}](#slide-{i}-{self._slugify(title)})")
        
        md_lines.extend(["", "---", "", "## Executive Summary", ""])
        
        # Add executive summary
        summary = self._generate_summary(prs)
        md_lines.append(summary)
        
        md_lines.extend(["", "---", "", "## Slides", ""])
        
        # Add each slide
        for i, slide in enumerate(prs.slides, 1):
            title = self._get_slide_title(slide) or f"Slide {i}"
            slug = self._slugify(title)
            
            md_lines.extend([
                f"### Slide {i}: {title}",
                f'<a name="slide-{i}-{slug}"></a>',
                ""
            ])
            
            # Add image if available
            if i <= len(images):
                image_path = images[i - 1]
                relative_path = f"images/{image_path.name}"
                md_lines.extend([
                    f"![Slide {i}]({relative_path})",
                    ""
                ])
            
            # Add content
            md_lines.extend([
                "**Content:**",
                ""
            ])
            
            content = self._get_slide_content(slide)
            if content:
                for line in content:
                    md_lines.append(f"- {line}")
            else:
                md_lines.append("*[No text content]*")
            
            md_lines.append("")
            
            # Add speaker notes
            notes = self._get_speaker_notes(slide)
            if notes:
                md_lines.extend([
                    "**Speaker Notes:**",
                    "",
                    f"> {notes}",
                    ""
                ])
            
            md_lines.extend(["", "---", ""])
        
        # Add appendix
        md_lines.extend([
            "## Appendix",
            "",
            "### Export Information",
            f"- Export Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            f"- Export Format: Markdown with PNG images",
            f"- Image Resolution: 1920x1080 (when available)",
            "",
            "### Resources",
            "- [Claude Code Documentation](https://docs.anthropic.com/claude-code)",
            "- [PowerPoint Best Practices](https://support.microsoft.com/office/powerpoint)",
            "",
            "### Contact",
            "Generated by Claude Code Slide Agent",
            "For issues or questions: https://github.com/schlessera/ppt-slide-agent",
            ""
        ])
        
        return "\n".join(md_lines)
    
    def _get_slide_title(self, slide) -> Optional[str]:
        """Extract title from a slide"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        
        # Look for the first text box
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if len(text) < 100:  # Likely a title
                    return text
                break
        
        return None
    
    def _get_slide_content(self, slide) -> List[str]:
        """Extract all text content from a slide"""
        content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Skip title if already extracted
                if shape == slide.shapes.title:
                    continue
                
                text = shape.text.strip()
                # Split by newlines and add each line
                for line in text.split('\n'):
                    if line.strip():
                        content.append(line.strip())
        
        return content
    
    def _get_speaker_notes(self, slide) -> Optional[str]:
        """Extract speaker notes from a slide"""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            return notes if notes else None
        return None
    
    def _generate_summary(self, prs: Presentation) -> str:
        """Generate an executive summary of the presentation"""
        titles = []
        for slide in prs.slides[:10]:  # Look at first 10 slides
            title = self._get_slide_title(slide)
            if title:
                titles.append(title)
        
        if not titles:
            return "This presentation contains visual content with minimal text."
        
        summary = f"This presentation consists of {len(prs.slides)} slides covering "
        
        if len(titles) > 3:
            summary += f"topics including {', '.join(titles[:3])}, and more."
        else:
            summary += f"{', '.join(titles)}."
        
        return summary
    
    def _slugify(self, text: str) -> str:
        """Convert text to URL-friendly slug"""
        import re
        text = text.lower()
        text = re.sub(r'[^\w\s-]', '', text)
        text = re.sub(r'[-\s]+', '-', text)
        return text.strip('-')[:50]
    
    def _generate_pdf(self, pptx_file: Path, output_dir: Path):
        """Generate PDF version of the presentation"""
        try:
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(pptx_file)
            ]
            
            subprocess.run(cmd, check=True)
            print(f"📑 Generated PDF: {output_dir / f'{pptx_file.stem}.pdf'}")
        except subprocess.CalledProcessError:
            print("⚠️  Could not generate PDF (LibreOffice may not be installed)")


def main():
    """Main entry point for the slide exporter"""
    import argparse
    
    parser = argparse.ArgumentParser(description="Export PowerPoint presentations to markdown")
    parser.add_argument("presentation", nargs="?", help="Path to PPTX file")
    parser.add_argument("--dir", default="./presentations", help="Presentations directory")
    parser.add_argument("--latest", action="store_true", help="Export the latest presentation")
    
    args = parser.parse_args()
    
    exporter = SlideExporter(args.dir)
    
    if args.latest:
        # Find the latest PPTX file
        pptx_files = list(Path(args.dir).glob("*.pptx"))
        if not pptx_files:
            print("❌ No presentations found")
            sys.exit(1)
        
        latest = max(pptx_files, key=lambda p: p.stat().st_mtime)
        print(f"📊 Exporting latest presentation: {latest}")
        exporter.export_to_markdown(latest)
    
    elif args.presentation:
        pptx_file = Path(args.presentation)
        if not pptx_file.exists():
            print(f"❌ File not found: {pptx_file}")
            sys.exit(1)
        
        exporter.export_to_markdown(pptx_file)
    
    else:
        print("❌ Please specify a presentation file or use --latest")
        sys.exit(1)


if __name__ == "__main__":
    main()