#!/usr/bin/env python3
"""
Design Applier Tool
Applies consistent design system to PowerPoint presentations
"""

import json
import logging
from pathlib import Path
from typing import Dict, List, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import colorsys

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class DesignSystem:
    """Design system configuration"""
    
    def __init__(self, config: Dict = None):
        """Initialize design system with configuration"""
        
        if config:
            self.config = config
        else:
            # Default design system
            self.config = {
                'colors': {
                    'primary': '#2563EB',
                    'primary_light': '#60A5FA',
                    'primary_dark': '#1E40AF',
                    'accent': '#F59E0B',
                    'accent_light': '#FCD34D',
                    'accent_dark': '#D97706',
                    'text_primary': '#1F2937',
                    'text_secondary': '#6B7280',
                    'text_tertiary': '#9CA3AF',
                    'background': '#FFFFFF',
                    'surface': '#F9FAFB',
                    'border': '#E5E7EB',
                    'success': '#10B981',
                    'warning': '#F59E0B',
                    'error': '#EF4444'
                },
                'typography': {
                    'font_family': 'Inter',
                    'font_family_fallback': 'Segoe UI',
                    'sizes': {
                        'hero': 76,
                        'h1': 57,
                        'h2': 43,
                        'h3': 32,
                        'body': 24,
                        'caption': 18,
                        'micro': 14
                    },
                    'weights': {
                        'black': 900,
                        'bold': 700,
                        'semibold': 600,
                        'regular': 400,
                        'light': 300
                    },
                    'line_heights': {
                        'tight': 1.1,
                        'normal': 1.5,
                        'relaxed': 1.75
                    }
                },
                'spacing': {
                    'base': 8,
                    'scale': [0, 4, 8, 12, 16, 24, 32, 48, 64, 96, 128]
                },
                'effects': {
                    'shadows': {
                        'sm': '0 1px 2px rgba(0,0,0,0.05)',
                        'md': '0 4px 6px rgba(0,0,0,0.07)',
                        'lg': '0 10px 15px rgba(0,0,0,0.10)'
                    },
                    'borders': {
                        'width': 1,
                        'radius': [0, 4, 8, 12, 16, 24]
                    }
                }
            }
    
    def get_color(self, color_name: str) -> RGBColor:
        """Get color as RGBColor object"""
        hex_color = self.config['colors'].get(color_name, '#000000')
        return self.hex_to_rgb(hex_color)
    
    @staticmethod
    def hex_to_rgb(hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    
    def get_font_size(self, size_name: str) -> Pt:
        """Get font size as Pt object"""
        size = self.config['typography']['sizes'].get(size_name, 24)
        return Pt(size)
    
    def get_spacing(self, level: int) -> float:
        """Get spacing value in inches"""
        if level < len(self.config['spacing']['scale']):
            pixels = self.config['spacing']['scale'][level]
            return pixels / 96.0  # Convert pixels to inches (96 DPI)
        return 0.5  # Default spacing


class DesignApplier:
    """Apply design system to presentations"""
    
    def __init__(self, presentation_path: Path, design_system: DesignSystem = None):
        """
        Initialize design applier
        
        Args:
            presentation_path: Path to PowerPoint file
            design_system: Design system to apply (uses default if None)
        """
        self.presentation_path = Path(presentation_path)
        self.presentation = Presentation(str(presentation_path))
        self.design_system = design_system or DesignSystem()
        
        logger.info(f"Loaded presentation: {presentation_path}")
    
    def apply_design_system(self):
        """Apply design system to entire presentation"""
        logger.info("Applying design system...")
        
        # Apply to each slide
        for i, slide in enumerate(self.presentation.slides):
            self.apply_slide_design(slide, i)
            logger.info(f"Applied design to slide {i+1}/{len(self.presentation.slides)}")
        
        # Save presentation
        output_path = self.presentation_path.parent / f"{self.presentation_path.stem}_styled{self.presentation_path.suffix}"
        self.presentation.save(str(output_path))
        
        logger.info(f"Styled presentation saved to: {output_path}")
        return output_path
    
    def apply_slide_design(self, slide, slide_index: int):
        """Apply design to individual slide"""
        
        # Determine slide type
        slide_type = self.detect_slide_type(slide)
        
        # Apply appropriate design based on type
        if slide_type == 'title':
            self.apply_title_slide_design(slide)
        elif slide_type == 'content':
            self.apply_content_slide_design(slide)
        elif slide_type == 'image':
            self.apply_image_slide_design(slide)
        else:
            self.apply_default_design(slide)
    
    def detect_slide_type(self, slide) -> str:
        """Detect the type of slide"""
        
        # Check if it's a title slide (first slide or has title layout)
        if slide.slide_layout.name in ['Title Slide', 'Title Only']:
            return 'title'
        
        # Check if it has large image
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                if shape.width > Inches(5):  # Large image
                    return 'image'
        
        # Default to content slide
        return 'content'
    
    def apply_title_slide_design(self, slide):
        """Apply design to title slide"""
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Apply text styling
                for paragraph in shape.text_frame.paragraphs:
                    # Check if it's the main title
                    if shape == slide.shapes.title:
                        self.style_text(
                            paragraph,
                            font_size=self.design_system.get_font_size('hero'),
                            font_color=self.design_system.get_color('primary_dark'),
                            bold=True,
                            alignment=PP_ALIGN.CENTER
                        )
                    else:
                        # Subtitle or other text
                        self.style_text(
                            paragraph,
                            font_size=self.design_system.get_font_size('h2'),
                            font_color=self.design_system.get_color('text_secondary'),
                            alignment=PP_ALIGN.CENTER
                        )
        
        # Apply background gradient (if supported)
        # Note: python-pptx has limited gradient support
        self.apply_slide_background(slide, 'gradient')
    
    def apply_content_slide_design(self, slide):
        """Apply design to content slide"""
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for i, paragraph in enumerate(shape.text_frame.paragraphs):
                    # Title
                    if shape == slide.shapes.title:
                        self.style_text(
                            paragraph,
                            font_size=self.design_system.get_font_size('h1'),
                            font_color=self.design_system.get_color('primary_dark'),
                            bold=True
                        )
                    # First level bullets
                    elif paragraph.level == 0:
                        self.style_text(
                            paragraph,
                            font_size=self.design_system.get_font_size('body'),
                            font_color=self.design_system.get_color('text_primary')
                        )
                    # Second level bullets
                    elif paragraph.level == 1:
                        self.style_text(
                            paragraph,
                            font_size=self.design_system.get_font_size('caption'),
                            font_color=self.design_system.get_color('text_secondary')
                        )
    
    def apply_image_slide_design(self, slide):
        """Apply design to image-heavy slide"""
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Text on image slides should be high contrast
                for paragraph in shape.text_frame.paragraphs:
                    if shape == slide.shapes.title:
                        self.style_text(
                            paragraph,
                            font_size=self.design_system.get_font_size('h1'),
                            font_color=self.design_system.get_color('primary_dark'),
                            bold=True
                        )
                    else:
                        self.style_text(
                            paragraph,
                            font_size=self.design_system.get_font_size('body'),
                            font_color=self.design_system.get_color('text_primary')
                        )
    
    def apply_default_design(self, slide):
        """Apply default design to slide"""
        self.apply_content_slide_design(slide)
    
    def style_text(self, paragraph, font_size=None, font_color=None, 
                   bold=None, italic=None, alignment=None):
        """Apply styling to text paragraph"""
        
        for run in paragraph.runs:
            if font_size:
                run.font.size = font_size
            if font_color:
                run.font.color.rgb = font_color
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
        
        if alignment:
            paragraph.alignment = alignment
        
        # Apply font family
        for run in paragraph.runs:
            run.font.name = self.design_system.config['typography']['font_family']
    
    def apply_slide_background(self, slide, background_type='solid'):
        """Apply background to slide"""
        
        # Note: python-pptx has limited background support
        # This is a simplified implementation
        
        if background_type == 'solid':
            # Would apply solid color background
            pass
        elif background_type == 'gradient':
            # Would apply gradient background
            pass
    
    def optimize_spacing(self, slide):
        """Optimize spacing between elements"""
        
        # Get all shapes with positions
        positioned_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                positioned_shapes.append(shape)
        
        # Sort by vertical position
        positioned_shapes.sort(key=lambda s: s.top)
        
        # Apply consistent spacing
        if len(positioned_shapes) > 1:
            spacing = self.design_system.get_spacing(6)  # Level 6 spacing
            current_top = positioned_shapes[0].top
            
            for shape in positioned_shapes[1:]:
                current_top += positioned_shapes[0].height + Inches(spacing)
                shape.top = current_top
    
    def apply_color_scheme(self, color_scheme: str = 'default'):
        """Apply a predefined color scheme"""
        
        schemes = {
            'default': {
                'primary': '#2563EB',
                'accent': '#F59E0B'
            },
            'corporate': {
                'primary': '#1E40AF',
                'accent': '#059669'
            },
            'creative': {
                'primary': '#DC2626',
                'accent': '#7C3AED'
            },
            'nature': {
                'primary': '#059669',
                'accent': '#F59E0B'
            }
        }
        
        if color_scheme in schemes:
            # Update design system colors
            self.design_system.config['colors'].update(schemes[color_scheme])
            logger.info(f"Applied {color_scheme} color scheme")
    
    def generate_style_report(self) -> Dict:
        """Generate report of applied styles"""
        
        report = {
            'total_slides': len(self.presentation.slides),
            'design_system': self.design_system.config,
            'slides_styled': [],
            'issues_found': [],
            'improvements_made': []
        }
        
        for i, slide in enumerate(self.presentation.slides):
            slide_info = {
                'slide_number': i + 1,
                'type': self.detect_slide_type(slide),
                'text_elements': 0,
                'images': 0,
                'shapes': len(slide.shapes)
            }
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_info['text_elements'] += 1
                if shape.shape_type == 13:  # Picture
                    slide_info['images'] += 1
            
            report['slides_styled'].append(slide_info)
        
        return report


def main():
    """Main function for command-line usage"""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python design_applier.py <presentation.pptx> [color_scheme]")
        print("Color schemes: default, corporate, creative, nature")
        sys.exit(1)
    
    presentation_path = Path(sys.argv[1])
    color_scheme = sys.argv[2] if len(sys.argv) > 2 else 'default'
    
    if not presentation_path.exists():
        print(f"Error: Presentation not found: {presentation_path}")
        sys.exit(1)
    
    # Apply design system
    applier = DesignApplier(presentation_path)
    
    # Apply color scheme if specified
    if color_scheme != 'default':
        applier.apply_color_scheme(color_scheme)
    
    # Apply design
    output_path = applier.apply_design_system()
    
    # Generate report
    report = applier.generate_style_report()
    
    print(f"\n✅ Design system applied successfully!")
    print(f"📁 Styled presentation: {output_path}")
    print(f"\n📊 Summary:")
    print(f"  • Slides styled: {report['total_slides']}")
    print(f"  • Color scheme: {color_scheme}")
    print(f"  • Primary color: {applier.design_system.config['colors']['primary']}")
    print(f"  • Accent color: {applier.design_system.config['colors']['accent']}")


if __name__ == "__main__":
    main()