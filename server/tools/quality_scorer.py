#!/usr/bin/env python3
"""
Quality Scorer Tool
Analyzes presentation quality across multiple dimensions
"""

import json
import logging
from pathlib import Path
from typing import Dict, List, Tuple, Optional
from dataclasses import dataclass, asdict
from pptx import Presentation
from pptx.util import Pt
import colorsys

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


@dataclass
class QualityScore:
    """Quality score for a single dimension"""
    dimension: str
    score: float  # 0-100
    weight: float  # 0-1
    issues: List[str]
    suggestions: List[str]


@dataclass
class SlideAnalysis:
    """Analysis results for a single slide"""
    slide_number: int
    word_count: int
    has_visual: bool
    font_sizes: List[float]
    contrast_issues: List[str]
    alignment_issues: List[str]
    overall_score: float


class QualityScorer:
    """Analyze and score presentation quality"""
    
    def __init__(self, presentation_path: Path):
        """
        Initialize the quality scorer
        
        Args:
            presentation_path: Path to the PowerPoint file
        """
        self.presentation_path = Path(presentation_path)
        self.presentation = Presentation(str(presentation_path))
        
        # Quality dimensions and weights
        self.dimensions = {
            'design_consistency': 0.20,
            'content_quality': 0.25,
            'visual_impact': 0.20,
            'accessibility': 0.15,
            'technical_quality': 0.10,
            'engagement_potential': 0.10
        }
    
    def analyze_presentation(self) -> Dict:
        """
        Perform complete quality analysis
        
        Returns:
            Dictionary containing all analysis results
        """
        results = {
            'overall_score': 0,
            'dimension_scores': {},
            'slide_analyses': [],
            'critical_issues': [],
            'recommendations': [],
            'summary': {}
        }
        
        # Analyze each dimension
        for dimension, weight in self.dimensions.items():
            score = self._analyze_dimension(dimension)
            results['dimension_scores'][dimension] = asdict(score)
            results['overall_score'] += score.score * weight
        
        # Analyze individual slides
        for i, slide in enumerate(self.presentation.slides):
            slide_analysis = self._analyze_slide(slide, i)
            results['slide_analyses'].append(asdict(slide_analysis))
        
        # Identify critical issues
        results['critical_issues'] = self._identify_critical_issues(results)
        
        # Generate recommendations
        results['recommendations'] = self._generate_recommendations(results)
        
        # Create summary
        results['summary'] = self._create_summary(results)
        
        return results
    
    def _analyze_dimension(self, dimension: str) -> QualityScore:
        """Analyze a specific quality dimension"""
        
        if dimension == 'design_consistency':
            return self._analyze_design_consistency()
        elif dimension == 'content_quality':
            return self._analyze_content_quality()
        elif dimension == 'visual_impact':
            return self._analyze_visual_impact()
        elif dimension == 'accessibility':
            return self._analyze_accessibility()
        elif dimension == 'technical_quality':
            return self._analyze_technical_quality()
        elif dimension == 'engagement_potential':
            return self._analyze_engagement_potential()
        else:
            return QualityScore(dimension, 0, 0, [], [])
    
    def _analyze_design_consistency(self) -> QualityScore:
        """Analyze design consistency across slides"""
        issues = []
        suggestions = []
        score = 100
        
        # Check font consistency
        fonts_used = set()
        font_sizes = set()
        
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                fonts_used.add(run.font.name)
                            if run.font.size:
                                font_sizes.add(run.font.size.pt)
        
        # Check for too many fonts
        if len(fonts_used) > 3:
            issues.append(f"Too many fonts used ({len(fonts_used)})")
            suggestions.append("Limit to 2-3 fonts maximum")
            score -= 15
        
        # Check for inconsistent font sizes
        if len(font_sizes) > 7:
            issues.append(f"Too many font size variations ({len(font_sizes)})")
            suggestions.append("Use a consistent type scale")
            score -= 10
        
        # Check color consistency (simplified)
        # In production, would analyze actual color usage
        
        return QualityScore(
            dimension='design_consistency',
            score=max(0, score),
            weight=self.dimensions['design_consistency'],
            issues=issues,
            suggestions=suggestions
        )
    
    def _analyze_content_quality(self) -> QualityScore:
        """Analyze content quality"""
        issues = []
        suggestions = []
        score = 100
        
        total_words = 0
        slides_over_limit = 0
        
        for i, slide in enumerate(self.presentation.slides):
            word_count = self._count_slide_words(slide)
            total_words += word_count
            
            if word_count > 40:
                slides_over_limit += 1
                issues.append(f"Slide {i+1}: {word_count} words (limit: 40)")
                score -= 5
        
        avg_words = total_words / len(self.presentation.slides) if self.presentation.slides else 0
        
        if avg_words > 35:
            suggestions.append(f"Reduce average word count from {avg_words:.1f} to <30")
            score -= 10
        
        if slides_over_limit > 0:
            suggestions.append(f"Simplify {slides_over_limit} text-heavy slides")
        
        return QualityScore(
            dimension='content_quality',
            score=max(0, score),
            weight=self.dimensions['content_quality'],
            issues=issues,
            suggestions=suggestions
        )
    
    def _analyze_visual_impact(self) -> QualityScore:
        """Analyze visual impact"""
        issues = []
        suggestions = []
        score = 100
        
        slides_without_visuals = 0
        
        for i, slide in enumerate(self.presentation.slides):
            if not self._has_visual_element(slide):
                slides_without_visuals += 1
                issues.append(f"Slide {i+1}: No visual element")
                score -= 3
        
        if slides_without_visuals > 0:
            suggestions.append(f"Add visuals to {slides_without_visuals} slides")
        
        # Check for visual hierarchy (simplified)
        # In production, would analyze actual visual weight distribution
        
        return QualityScore(
            dimension='visual_impact',
            score=max(0, score),
            weight=self.dimensions['visual_impact'],
            issues=issues,
            suggestions=suggestions
        )
    
    def _analyze_accessibility(self) -> QualityScore:
        """Analyze accessibility compliance"""
        issues = []
        suggestions = []
        score = 100
        
        # Check font sizes
        small_font_count = 0
        
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size.pt < 24:
                                small_font_count += 1
        
        if small_font_count > 0:
            issues.append(f"{small_font_count} text elements below 24pt")
            suggestions.append("Increase all body text to minimum 24pt")
            score -= min(30, small_font_count * 2)
        
        # Check for alt text (simplified - would need to check images)
        # Check color contrast (simplified - would need color analysis)
        
        return QualityScore(
            dimension='accessibility',
            score=max(0, score),
            weight=self.dimensions['accessibility'],
            issues=issues,
            suggestions=suggestions
        )
    
    def _analyze_technical_quality(self) -> QualityScore:
        """Analyze technical quality"""
        issues = []
        suggestions = []
        score = 100
        
        # Check presentation size (simplified)
        file_size_mb = self.presentation_path.stat().st_size / (1024 * 1024)
        
        if file_size_mb > 50:
            issues.append(f"File size too large: {file_size_mb:.1f}MB")
            suggestions.append("Compress images and remove unused elements")
            score -= 20
        elif file_size_mb > 25:
            issues.append(f"File size could be optimized: {file_size_mb:.1f}MB")
            suggestions.append("Consider compressing images")
            score -= 10
        
        return QualityScore(
            dimension='technical_quality',
            score=max(0, score),
            weight=self.dimensions['technical_quality'],
            issues=issues,
            suggestions=suggestions
        )
    
    def _analyze_engagement_potential(self) -> QualityScore:
        """Analyze engagement potential"""
        issues = []
        suggestions = []
        score = 85  # Start with good baseline
        
        # Check for speaker notes
        slides_with_notes = sum(1 for slide in self.presentation.slides if slide.has_notes_slide)
        
        if slides_with_notes < len(self.presentation.slides) * 0.8:
            issues.append(f"Only {slides_with_notes} slides have speaker notes")
            suggestions.append("Add comprehensive speaker notes to all slides")
            score -= 15
        
        # Check for variety in slide layouts (simplified)
        
        return QualityScore(
            dimension='engagement_potential',
            score=max(0, score),
            weight=self.dimensions['engagement_potential'],
            issues=issues,
            suggestions=suggestions
        )
    
    def _analyze_slide(self, slide, index: int) -> SlideAnalysis:
        """Analyze individual slide"""
        word_count = self._count_slide_words(slide)
        has_visual = self._has_visual_element(slide)
        font_sizes = self._get_font_sizes(slide)
        
        # Calculate slide score
        slide_score = 100
        if word_count > 40:
            slide_score -= 20
        if not has_visual:
            slide_score -= 15
        if min(font_sizes) < 24 if font_sizes else False:
            slide_score -= 10
        
        return SlideAnalysis(
            slide_number=index + 1,
            word_count=word_count,
            has_visual=has_visual,
            font_sizes=font_sizes,
            contrast_issues=[],  # Simplified
            alignment_issues=[],  # Simplified
            overall_score=max(0, slide_score)
        )
    
    def _count_slide_words(self, slide) -> int:
        """Count words on a slide"""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_parts.append(shape.text)
        
        all_text = ' '.join(text_parts)
        return len(all_text.split())
    
    def _has_visual_element(self, slide) -> bool:
        """Check if slide has visual elements"""
        for shape in slide.shapes:
            # Check for pictures
            if shape.shape_type == 13:  # Picture
                return True
            # Check for charts
            if hasattr(shape, 'has_chart') and shape.has_chart:
                return True
            # Check for tables
            if hasattr(shape, 'has_table') and shape.has_table:
                return True
        return False
    
    def _get_font_sizes(self, slide) -> List[float]:
        """Get all font sizes used in slide"""
        font_sizes = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)
        
        return font_sizes
    
    def _identify_critical_issues(self, results: Dict) -> List[str]:
        """Identify critical issues that must be fixed"""
        critical = []
        
        # Check for accessibility failures
        if results['dimension_scores']['accessibility']['score'] < 70:
            critical.append("Accessibility score below acceptable threshold")
        
        # Check for severe text overload
        severely_overloaded = sum(
            1 for s in results['slide_analyses'] 
            if s['word_count'] > 60
        )
        if severely_overloaded > 0:
            critical.append(f"{severely_overloaded} slides severely text-heavy (>60 words)")
        
        return critical
    
    def _generate_recommendations(self, results: Dict) -> List[Dict]:
        """Generate prioritized recommendations"""
        recommendations = []
        
        # Collect all suggestions
        for dimension_data in results['dimension_scores'].values():
            for suggestion in dimension_data['suggestions']:
                recommendations.append({
                    'dimension': dimension_data['dimension'],
                    'suggestion': suggestion,
                    'priority': 'high' if dimension_data['score'] < 70 else 'medium'
                })
        
        # Sort by priority
        priority_order = {'critical': 0, 'high': 1, 'medium': 2, 'low': 3}
        recommendations.sort(key=lambda x: priority_order.get(x['priority'], 3))
        
        return recommendations[:10]  # Top 10 recommendations
    
    def _create_summary(self, results: Dict) -> Dict:
        """Create executive summary"""
        total_slides = len(self.presentation.slides)
        avg_score = sum(s['overall_score'] for s in results['slide_analyses']) / total_slides if total_slides else 0
        
        return {
            'total_slides': total_slides,
            'overall_score': round(results['overall_score'], 1),
            'average_slide_score': round(avg_score, 1),
            'critical_issues_count': len(results['critical_issues']),
            'status': self._get_status(results['overall_score']),
            'top_dimension': max(
                results['dimension_scores'].items(),
                key=lambda x: x[1]['score']
            )[0],
            'weakest_dimension': min(
                results['dimension_scores'].items(),
                key=lambda x: x[1]['score']
            )[0]
        }
    
    def _get_status(self, score: float) -> str:
        """Get status based on score"""
        if score >= 90:
            return "Excellent - Ready for presentation"
        elif score >= 80:
            return "Good - Minor improvements recommended"
        elif score >= 70:
            return "Satisfactory - Several improvements needed"
        elif score >= 60:
            return "Poor - Significant improvements required"
        else:
            return "Critical - Major revision needed"
    
    def generate_report(self, output_path: Path = None) -> Path:
        """
        Generate quality report
        
        Args:
            output_path: Path to save report (default: quality_report.json)
            
        Returns:
            Path to generated report
        """
        results = self.analyze_presentation()
        
        output_path = output_path or Path('quality_report.json')
        
        with open(output_path, 'w') as f:
            json.dump(results, f, indent=2)
        
        logger.info(f"Quality report saved to: {output_path}")
        
        # Print summary to console
        print(f"\n{'='*50}")
        print(f"PRESENTATION QUALITY REPORT")
        print(f"{'='*50}")
        print(f"Overall Score: {results['overall_score']:.1f}/100")
        print(f"Status: {results['summary']['status']}")
        print(f"\nDimension Scores:")
        for dimension, data in results['dimension_scores'].items():
            print(f"  • {dimension}: {data['score']:.1f}/100")
        
        if results['critical_issues']:
            print(f"\n⚠️  Critical Issues:")
            for issue in results['critical_issues']:
                print(f"  • {issue}")
        
        print(f"\n📋 Top Recommendations:")
        for i, rec in enumerate(results['recommendations'][:5], 1):
            print(f"  {i}. {rec['suggestion']}")
        
        return output_path


def main():
    """Main function for command-line usage"""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python quality_scorer.py <presentation.pptx> [output.json]")
        sys.exit(1)
    
    presentation_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else None
    
    if not presentation_path.exists():
        print(f"Error: Presentation not found: {presentation_path}")
        sys.exit(1)
    
    # Analyze presentation
    scorer = QualityScorer(presentation_path)
    report_path = scorer.generate_report(output_path)
    
    print(f"\n✅ Analysis complete!")
    print(f"📊 Full report saved to: {report_path}")


if __name__ == "__main__":
    main()