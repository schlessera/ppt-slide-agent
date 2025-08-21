#!/usr/bin/env python3
"""
Screenshot Extractor Tool
Extracts PNG screenshots from PowerPoint presentations for quality review
"""

import os
import sys
from pathlib import Path
from typing import List, Tuple, Optional
import logging
from pptx import Presentation
from PIL import Image
import io
import base64

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class ScreenshotExtractor:
    """Extract screenshots from PowerPoint presentations"""
    
    def __init__(self, presentation_path: Path, output_dir: Path = None):
        """
        Initialize the screenshot extractor
        
        Args:
            presentation_path: Path to the PowerPoint file
            output_dir: Directory to save screenshots (default: ./screenshots)
        """
        self.presentation_path = Path(presentation_path)
        self.output_dir = output_dir or Path('./screenshots')
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Load presentation
        try:
            self.presentation = Presentation(str(self.presentation_path))
            logger.info(f"Loaded presentation: {self.presentation_path}")
        except Exception as e:
            logger.error(f"Failed to load presentation: {e}")
            raise
    
    def extract_all_slides(self) -> List[Path]:
        """
        Extract screenshots of all slides
        
        Returns:
            List of paths to generated screenshot files
        """
        screenshot_paths = []
        
        for i, slide in enumerate(self.presentation.slides):
            try:
                screenshot_path = self.extract_slide(i)
                screenshot_paths.append(screenshot_path)
                logger.info(f"Extracted slide {i+1}/{len(self.presentation.slides)}")
            except Exception as e:
                logger.error(f"Failed to extract slide {i+1}: {e}")
        
        return screenshot_paths
    
    def extract_slide(self, slide_index: int) -> Path:
        """
        Extract a single slide as PNG
        
        Args:
            slide_index: Index of the slide to extract
            
        Returns:
            Path to the generated screenshot
        """
        slide = self.presentation.slides[slide_index]
        
        # Generate image (this is a simplified version - in production, 
        # you'd use win32com on Windows or python-pptx-interface on Linux)
        image = self._render_slide_to_image(slide)
        
        # Save image
        output_path = self.output_dir / f"slide_{slide_index+1:02d}.png"
        image.save(output_path, 'PNG', optimize=True)
        
        return output_path
    
    def _render_slide_to_image(self, slide) -> Image:
        """
        Render slide to image (simplified implementation)
        
        In production, this would use:
        - win32com.client on Windows to control PowerPoint
        - LibreOffice headless on Linux
        - Or a cloud service API
        """
        # For now, create a placeholder image with slide information
        # In production, this would actually render the slide
        
        img = Image.new('RGB', (1920, 1080), color='white')
        
        # This is where you'd implement actual slide rendering
        # Options:
        # 1. Use win32com to control PowerPoint (Windows only)
        # 2. Use LibreOffice in headless mode (cross-platform)
        # 3. Use a cloud rendering service
        # 4. Use python-pptx-interface for basic rendering
        
        return img
    
    def extract_with_metadata(self) -> List[dict]:
        """
        Extract slides with metadata for quality review
        
        Returns:
            List of dictionaries containing screenshot paths and metadata
        """
        results = []
        
        for i, slide in enumerate(self.presentation.slides):
            # Extract text content
            text_content = self._extract_slide_text(slide)
            word_count = len(text_content.split())
            
            # Extract screenshot
            screenshot_path = self.extract_slide(i)
            
            # Gather metadata
            metadata = {
                'slide_number': i + 1,
                'screenshot_path': str(screenshot_path),
                'word_count': word_count,
                'has_title': bool(slide.shapes.title),
                'shape_count': len(slide.shapes),
                'has_image': self._has_image(slide),
                'has_chart': self._has_chart(slide),
                'text_content': text_content[:200],  # First 200 chars
                'notes': self._get_speaker_notes(slide)
            }
            
            results.append(metadata)
        
        return results
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a slide"""
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_parts.append(shape.text)
        
        return ' '.join(text_parts)
    
    def _has_image(self, slide) -> bool:
        """Check if slide contains an image"""
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                return True
        return False
    
    def _has_chart(self, slide) -> bool:
        """Check if slide contains a chart"""
        for shape in slide.shapes:
            if shape.has_chart:
                return True
        return False
    
    def _get_speaker_notes(self, slide) -> str:
        """Extract speaker notes from slide"""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            return notes_text
        return ""
    
    def generate_review_html(self, metadata: List[dict]) -> Path:
        """
        Generate an HTML review page with all screenshots
        
        Args:
            metadata: List of slide metadata dictionaries
            
        Returns:
            Path to generated HTML file
        """
        html_content = """
        <!DOCTYPE html>
        <html>
        <head>
            <title>Presentation Review</title>
            <style>
                body { font-family: Arial, sans-serif; margin: 20px; }
                .slide { margin-bottom: 40px; border: 1px solid #ddd; padding: 20px; }
                .screenshot { max-width: 800px; border: 1px solid #ccc; }
                .metadata { margin-top: 10px; }
                .warning { color: orange; }
                .error { color: red; }
                .success { color: green; }
            </style>
        </head>
        <body>
            <h1>Presentation Quality Review</h1>
        """
        
        for slide_data in metadata:
            status_class = 'success'
            if slide_data['word_count'] > 40:
                status_class = 'error'
            elif slide_data['word_count'] > 30:
                status_class = 'warning'
            
            html_content += f"""
            <div class="slide">
                <h2>Slide {slide_data['slide_number']}</h2>
                <img src="{slide_data['screenshot_path']}" class="screenshot" />
                <div class="metadata">
                    <p>Word Count: <span class="{status_class}">{slide_data['word_count']}</span></p>
                    <p>Has Image: {'✓' if slide_data['has_image'] else '✗'}</p>
                    <p>Has Chart: {'✓' if slide_data['has_chart'] else '✗'}</p>
                    <p>Shape Count: {slide_data['shape_count']}</p>
                </div>
            </div>
            """
        
        html_content += """
        </body>
        </html>
        """
        
        # Save HTML
        html_path = self.output_dir / 'review.html'
        html_path.write_text(html_content)
        
        return html_path


def main():
    """Main function for command-line usage"""
    if len(sys.argv) < 2:
        print("Usage: python screenshot_extractor.py <presentation.pptx> [output_dir]")
        sys.exit(1)
    
    presentation_path = Path(sys.argv[1])
    output_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else None
    
    if not presentation_path.exists():
        print(f"Error: Presentation file not found: {presentation_path}")
        sys.exit(1)
    
    # Extract screenshots
    extractor = ScreenshotExtractor(presentation_path, output_dir)
    
    # Extract with metadata
    metadata = extractor.extract_with_metadata()
    
    # Generate review HTML
    html_path = extractor.generate_review_html(metadata)
    
    print(f"✅ Extracted {len(metadata)} slides")
    print(f"📁 Screenshots saved to: {extractor.output_dir}")
    print(f"🌐 Review page: {html_path}")
    
    # Print summary
    total_words = sum(s['word_count'] for s in metadata)
    avg_words = total_words / len(metadata) if metadata else 0
    
    print(f"\n📊 Summary:")
    print(f"  • Total slides: {len(metadata)}")
    print(f"  • Average words per slide: {avg_words:.1f}")
    print(f"  • Slides with images: {sum(1 for s in metadata if s['has_image'])}")
    print(f"  • Slides over 40 words: {sum(1 for s in metadata if s['word_count'] > 40)}")


if __name__ == "__main__":
    main()