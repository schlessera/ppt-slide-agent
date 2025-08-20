#!/usr/bin/env python3
"""
PPTX MCP Server - FastMCP-powered PowerPoint creation server
Based on https://github.com/samos123/pptx-mcp
Enhanced with export capabilities for Claude Code Slide Agent
"""

import os
import sys
import json
import base64
import subprocess
from pathlib import Path
from typing import Optional, List, Dict, Any
from datetime import datetime

from fastmcp import FastMCP, Context
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import io

# Initialize FastMCP server
mcp = FastMCP("pptx-mcp-server")

# Configuration
HOST = os.getenv("HOST", "127.0.0.1")
PORT = int(os.getenv("PORT", "8000"))
PRESENTATIONS_DIR = Path(os.getenv("PRESENTATIONS_DIR", "./presentations"))
TEMPLATES_DIR = PRESENTATIONS_DIR / "templates"
EXPORTS_DIR = PRESENTATIONS_DIR / "exports"

# Ensure directories exist
PRESENTATIONS_DIR.mkdir(parents=True, exist_ok=True)
TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
EXPORTS_DIR.mkdir(parents=True, exist_ok=True)

# Store active presentations in memory
presentations: Dict[str, Presentation] = {}

@mcp.tool()
def create_presentation(name: str, template: Optional[str] = None) -> str:
    """
    Create a new PowerPoint presentation
    
    Args:
        name: Name for the presentation
        template: Optional template name to use
    
    Returns:
        Success message
    """
    try:
        if template and (TEMPLATES_DIR / f"{template}.pptx").exists():
            prs = Presentation(TEMPLATES_DIR / f"{template}.pptx")
        else:
            prs = Presentation()
        
        presentations[name] = prs
        return f"Created presentation '{name}'"
    except Exception as e:
        return f"Error creating presentation: {str(e)}"

@mcp.tool()
def add_slide(presentation_name: str, layout: str = "title_and_content") -> str:
    """
    Add a new slide to the presentation
    
    Args:
        presentation_name: Name of the presentation
        layout: Slide layout type (title_slide, title_and_content, blank, etc.)
    
    Returns:
        Slide index
    """
    if presentation_name not in presentations:
        return f"Presentation '{presentation_name}' not found"
    
    prs = presentations[presentation_name]
    
    # Map layout names to indices
    layout_map = {
        "title_slide": 0,
        "title_and_content": 1,
        "section_header": 2,
        "two_content": 3,
        "comparison": 4,
        "title_only": 5,
        "blank": 6,
        "content_with_caption": 7,
        "picture_with_caption": 8
    }
    
    layout_idx = layout_map.get(layout, 1)
    slide_layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(slide_layout)
    
    return f"Added slide {len(prs.slides) - 1} with layout '{layout}'"

@mcp.tool()
def add_text_to_slide(
    presentation_name: str,
    slide_index: int,
    text: str,
    placeholder_index: int = 0,
    font_size: Optional[int] = None,
    bold: bool = False,
    italic: bool = False,
    color: Optional[str] = None
) -> str:
    """
    Add text to a slide
    
    Args:
        presentation_name: Name of the presentation
        slide_index: Index of the slide
        text: Text to add
        placeholder_index: Index of the placeholder (0 for title, 1 for content)
        font_size: Optional font size in points
        bold: Make text bold
        italic: Make text italic
        color: Optional hex color (e.g., "#FF0000")
    
    Returns:
        Success message
    """
    if presentation_name not in presentations:
        return f"Presentation '{presentation_name}' not found"
    
    prs = presentations[presentation_name]
    
    if slide_index >= len(prs.slides):
        return f"Slide index {slide_index} out of range"
    
    slide = prs.slides[slide_index]
    
    # Try to find placeholder
    if hasattr(slide, 'placeholders') and len(slide.placeholders) > placeholder_index:
        placeholder = slide.placeholders[placeholder_index]
        placeholder.text = text
        
        # Apply formatting if specified
        if any([font_size, bold, italic, color]):
            for paragraph in placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    if font_size:
                        run.font.size = Pt(font_size)
                    if bold:
                        run.font.bold = True
                    if italic:
                        run.font.italic = True
                    if color:
                        rgb = tuple(int(color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
                        run.font.color.rgb = RGBColor(*rgb)
    else:
        # Add textbox if no placeholder
        left = Inches(1)
        top = Inches(1) if placeholder_index == 0 else Inches(2)
        width = Inches(8)
        height = Inches(1)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        
        if font_size:
            text_frame.paragraphs[0].font.size = Pt(font_size)
    
    return f"Added text to slide {slide_index}"

@mcp.tool()
def add_image_to_slide(
    presentation_name: str,
    slide_index: int,
    image_path: str,
    left: float = 1,
    top: float = 2,
    width: Optional[float] = None,
    height: Optional[float] = None
) -> str:
    """
    Add an image to a slide
    
    Args:
        presentation_name: Name of the presentation
        slide_index: Index of the slide
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Optional width in inches
        height: Optional height in inches
    
    Returns:
        Success message
    """
    if presentation_name not in presentations:
        return f"Presentation '{presentation_name}' not found"
    
    prs = presentations[presentation_name]
    
    if slide_index >= len(prs.slides):
        return f"Slide index {slide_index} out of range"
    
    if not Path(image_path).exists():
        return f"Image file '{image_path}' not found"
    
    slide = prs.slides[slide_index]
    
    left_pos = Inches(left)
    top_pos = Inches(top)
    
    if width and height:
        slide.shapes.add_picture(image_path, left_pos, top_pos, 
                                  width=Inches(width), height=Inches(height))
    elif width:
        slide.shapes.add_picture(image_path, left_pos, top_pos, width=Inches(width))
    elif height:
        slide.shapes.add_picture(image_path, left_pos, top_pos, height=Inches(height))
    else:
        slide.shapes.add_picture(image_path, left_pos, top_pos)
    
    return f"Added image to slide {slide_index}"

@mcp.tool()
def add_speaker_notes(presentation_name: str, slide_index: int, notes: str) -> str:
    """
    Add speaker notes to a slide
    
    Args:
        presentation_name: Name of the presentation
        slide_index: Index of the slide
        notes: Speaker notes text
    
    Returns:
        Success message
    """
    if presentation_name not in presentations:
        return f"Presentation '{presentation_name}' not found"
    
    prs = presentations[presentation_name]
    
    if slide_index >= len(prs.slides):
        return f"Slide index {slide_index} out of range"
    
    slide = prs.slides[slide_index]
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    
    return f"Added speaker notes to slide {slide_index}"

@mcp.tool()
def get_slide_content(presentation_name: str, slide_index: int) -> Dict[str, Any]:
    """
    Get the content of a specific slide
    
    Args:
        presentation_name: Name of the presentation
        slide_index: Index of the slide
    
    Returns:
        Dictionary with slide content
    """
    if presentation_name not in presentations:
        return {"error": f"Presentation '{presentation_name}' not found"}
    
    prs = presentations[presentation_name]
    
    if slide_index >= len(prs.slides):
        return {"error": f"Slide index {slide_index} out of range"}
    
    slide = prs.slides[slide_index]
    content = {
        "slide_index": slide_index,
        "layout": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "unknown",
        "text_content": [],
        "shapes": [],
        "notes": ""
    }
    
    # Extract text content
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            content["text_content"].append(shape.text)
        content["shapes"].append(shape.shape_type)
    
    # Extract speaker notes
    if slide.has_notes_slide:
        content["notes"] = slide.notes_slide.notes_text_frame.text
    
    return content

@mcp.tool()
def download_presentation(presentation_name: str) -> Dict[str, str]:
    """
    Save and download the presentation
    
    Args:
        presentation_name: Name of the presentation
    
    Returns:
        Dictionary with file path and status
    """
    if presentation_name not in presentations:
        return {"error": f"Presentation '{presentation_name}' not found"}
    
    prs = presentations[presentation_name]
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{presentation_name}_{timestamp}.pptx"
    filepath = PRESENTATIONS_DIR / filename
    
    try:
        prs.save(filepath)
        return {
            "status": "success",
            "filepath": str(filepath),
            "filename": filename,
            "message": f"Presentation saved to {filepath}"
        }
    except Exception as e:
        return {"error": f"Failed to save presentation: {str(e)}"}

@mcp.tool()
def list_presentations() -> List[str]:
    """
    List all active presentations in memory
    
    Returns:
        List of presentation names
    """
    return list(presentations.keys())

@mcp.tool()
def get_presentation_info(presentation_name: str) -> Dict[str, Any]:
    """
    Get information about a presentation
    
    Args:
        presentation_name: Name of the presentation
    
    Returns:
        Dictionary with presentation information
    """
    if presentation_name not in presentations:
        return {"error": f"Presentation '{presentation_name}' not found"}
    
    prs = presentations[presentation_name]
    
    info = {
        "name": presentation_name,
        "slide_count": len(prs.slides),
        "slides": []
    }
    
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i,
            "layout": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "unknown",
            "has_notes": slide.has_notes_slide
        }
        
        # Get title if available
        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text
        
        info["slides"].append(slide_info)
    
    return info

@mcp.tool()
def clear_presentation(presentation_name: str) -> str:
    """
    Clear/delete a presentation from memory
    
    Args:
        presentation_name: Name of the presentation
    
    Returns:
        Success message
    """
    if presentation_name in presentations:
        del presentations[presentation_name]
        return f"Cleared presentation '{presentation_name}'"
    return f"Presentation '{presentation_name}' not found"

@mcp.tool()
def render_slide_to_image(presentation_name: str, slide_index: int) -> Dict[str, Any]:
    """
    Render a slide to a PNG image using LibreOffice
    
    Args:
        presentation_name: Name of the presentation
        slide_index: Index of the slide to render
    
    Returns:
        Dictionary with image path or error
    """
    if presentation_name not in presentations:
        return {"error": f"Presentation '{presentation_name}' not found"}
    
    prs = presentations[presentation_name]
    
    if slide_index >= len(prs.slides):
        return {"error": f"Slide index {slide_index} out of range"}
    
    # Save presentation temporarily
    temp_path = PRESENTATIONS_DIR / f"temp_{presentation_name}.pptx"
    prs.save(temp_path)
    
    # Use LibreOffice to convert to images
    try:
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "png",
            "--outdir",
            str(EXPORTS_DIR),
            str(temp_path)
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        if result.returncode == 0:
            # Find the generated image
            image_pattern = temp_path.stem
            images = list(EXPORTS_DIR.glob(f"{image_pattern}*.png"))
            
            if images and slide_index < len(images):
                image_path = images[slide_index]
                return {
                    "status": "success",
                    "image_path": str(image_path),
                    "message": f"Slide {slide_index} rendered to {image_path}"
                }
        
        return {"error": f"Failed to render slide: {result.stderr}"}
    
    except Exception as e:
        return {"error": f"Error rendering slide: {str(e)}"}
    
    finally:
        # Clean up temp file
        if temp_path.exists():
            temp_path.unlink()

if __name__ == "__main__":
    print(f"Starting PPTX MCP Server on {HOST}:{PORT}")
    print(f"Presentations directory: {PRESENTATIONS_DIR}")
    print(f"Templates directory: {TEMPLATES_DIR}")
    print(f"Exports directory: {EXPORTS_DIR}")
    
    # Run the server
    import uvicorn
    uvicorn.run(mcp, host=HOST, port=PORT)